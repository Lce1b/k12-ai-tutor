"""
Resources Agent — generates PPT and Word documents for teaching.
"""

import io
import json
import os
import tempfile
from openai import AsyncOpenAI

from config import LLM_API_KEY, LLM_BASE_URL, LLM_MODEL
from agents.orchestrator import GradeLevel

client = AsyncOpenAI(api_key=LLM_API_KEY, base_url=LLM_BASE_URL)

PPTX_AVAILABLE = False
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    pass

DOCX_AVAILABLE = False
try:
    from docx import Document
    from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor
    DOCX_AVAILABLE = True
except ImportError:
    pass


async def generate_ppt(topic: str, grade: GradeLevel) -> bytes:
    """Generate a teaching PPT about a topic, returns .pptx bytes."""
    # Step 1: Use LLM to generate slide content
    prompt = f"""为「{topic}」生成一份K12教学PPT内容（{grade.value}水平）。

返回JSON：
{{
  "title": "PPT标题",
  "slides": [
    {{"heading": "标题页", "bullets": ["要点1", "要点2", "要点3"]}},
    {{"heading": "内容页标题", "bullets": ["要点1", "要点2"], "note": "讲解备注"}}
  ]
}}

生成6-8页幻灯片。bullets每页2-4条。"""

    resp = await client.chat.completions.create(
        model=LLM_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=2000,
        response_format={"type": "json_object"},
    )

    try:
        content = json.loads(resp.choices[0].message.content)
    except json.JSONDecodeError:
        content = {"title": topic, "slides": [{"heading": topic, "bullets": ["内容加载中..."]}]}

    if not PPTX_AVAILABLE:
        return _make_simple_pptx(content)

    return _make_pptx(content)


def _make_simple_pptx(content: dict) -> bytes:
    """Generate PPTX without python-pptx (simple placeholder)."""
    # Build a minimal valid PPTX (zip of XML files)
    return _build_minimal_pptx_zip(content)


def _make_pptx(content: dict) -> bytes:
    """Generate PPTX using python-pptx."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else slide.shapes.title
    if title:
        title.text = content.get("title", "教学课件")
    if subtitle and subtitle != title:
        subtitle.text = "K12 AI 智能教学助手"

    # Content slides
    for s in content.get("slides", []):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = s.get("heading", "")
        body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body:
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(s.get("bullets", [])):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _build_minimal_pptx_zip(content: dict) -> bytes:
    """Build a minimal valid PPTX from scratch (no python-pptx dependency)."""
    import zipfile

    buf = io.BytesIO()
    slides_xml = []
    slide_count = 1  # title slide + content slides

    # Build slide XMLs
    title = content.get("title", "教学课件")
    slides_data = content.get("slides", [])

    # Title slide
    slides_xml.append(f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:cSld>
    <p:spTree>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="1" name="Title"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="914400" y="2743200"/><a:ext cx="11430000" cy="1473200"/></a:xfrm></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="zh-CN" sz="4400" b="1"/><a:t>{_xml_escape(title)}</a:t></a:r></a:p></p:txBody>
      </p:sp>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Subtitle"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="914400" y="4572000"/><a:ext cx="11430000" cy="914400"/></a:xfrm></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="zh-CN" sz="2400"/><a:t>K12 AI 智能教学助手</a:t></a:r></a:p></p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
</p:sld>''')

    for i, s in enumerate(slides_data):
        heading = _xml_escape(s.get("heading", ""))
        bullets = "".join(
            f'<a:p><a:r><a:rPr lang="zh-CN" sz="2400"/><a:t>{_xml_escape(b)}</a:t></a:r></a:p>'
            for b in s.get("bullets", [])
        )
        slides_xml.append(f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:cSld>
    <p:spTree>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="1" name="Title"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="685800" y="342900"/><a:ext cx="12000000" cy="914400"/></a:xfrm></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="zh-CN" sz="3200" b="1"/><a:t>{heading}</a:t></a:r></a:p></p:txBody>
      </p:sp>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Body"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="914400" y="1371600"/><a:ext cx="11430000" cy="5200000"/></a:xfrm></p:spPr>
        <p:txBody><a:bodyPr/><a:lstStyle/>{bullets}</p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
</p:sld>''')
        slide_count += 1

    slide_refs = "".join(
        f'<p:sldId id="{256 + i}" r:id="rId{i + 1}"/>'
        for i in range(slide_count)
    )

    # Build the zip
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
  <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
  <Override PartName="/ppt/slideLayouts/slideLayout2.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
  <Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
  ''' + "".join(
      f'<Override PartName="/ppt/slides/slide{i + 1}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
      for i in range(slide_count)
  ) + "</Types>")
        z.writestr("_rels/.rels", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>''')
        z.writestr("ppt/_rels/presentation.xml.rels", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="theme/theme1.xml"/>
  ''' + "".join(
      f'<Relationship Id="rId{i + 3}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{i + 1}.xml"/>'
      for i in range(slide_count)
  ) + "</Relationships>")
        z.writestr("ppt/presentation.xml", f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
  <p:sldIdLst>{slide_refs}</p:sldIdLst>
  <p:sldSz cx="12192000" cy="6858000"/>
</p:presentation>''')
        z.writestr("ppt/slideMasters/slideMaster1.xml", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
             xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:cSld><p:spTree><p:sp><p:nvSpPr><p:cNvPr id="1" name=""/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="12192000" cy="6858000"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/></p:txBody></p:sp></p:spTree></p:cSld>
  <p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/><p:sldLayoutId id="2147483650" r:id="rId2"/></p:sldLayoutIdLst>
</p:sldMaster>''')
        z.writestr("ppt/slideLayouts/slideLayout1.xml", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld name="Title Slide"><p:spTree><p:sp><p:nvSpPr><p:cNvPr id="1" name="Title"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="12192000" cy="6858000"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/></p:txBody></p:sp>
  <p:sp><p:nvSpPr><p:cNvPr id="2" name="Subtitle"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="12192000" cy="6858000"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/></p:txBody></p:sp></p:spTree></p:cSld>
</p:sldLayout>''')
        z.writestr("ppt/slideLayouts/slideLayout2.xml", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld name="Content Slide"><p:spTree><p:sp><p:nvSpPr><p:cNvPr id="1" name="Title"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="12192000" cy="6858000"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/></p:txBody></p:sp>
  <p:sp><p:nvSpPr><p:cNvPr id="2" name="Body"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="12192000" cy="6858000"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/></p:txBody></p:sp></p:spTree></p:cSld>
</p:sldLayout>''')
        # Minimal theme
        z.writestr("ppt/theme/theme1.xml", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Default">
  <a:themeElements>
    <a:clrScheme name="Default"><a:dk1><a:srgbClr val="000000"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1>
    <a:dk2><a:srgbClr val="44546A"/></a:dk2><a:lt2><a:srgbClr val="E7E6E6"/></a:lt2>
    <a:accent1><a:srgbClr val="4472C4"/></a:accent1><a:accent2><a:srgbClr val="ED7D31"/></a:accent2>
    <a:accent3><a:srgbClr val="A5A5A5"/></a:accent3><a:accent4><a:srgbClr val="FFC000"/></a:accent4>
    <a:accent5><a:srgbClr val="5B9BD5"/></a:accent5><a:accent6><a:srgbClr val="70AD47"/></a:accent6>
    </a:clrScheme>
    <a:fontScheme name="Default"><a:majorFont><a:latin typeface="Calibri"/><a:ea typeface="Microsoft YaHei"/></a:majorFont>
    <a:minorFont><a:latin typeface="Calibri"/><a:ea typeface="Microsoft YaHei"/></a:minorFont></a:fontScheme>
    <a:fmtScheme name="Default"><a:fillStyleLst><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:fillStyleLst>
    <a:lnStyleLst><a:ln w="6350"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:ln></a:lnStyleLst>
    </a:fmtScheme>
  </a:themeElements>
</a:theme>''')
        for i, xml in enumerate(slides_xml):
            z.writestr(f"ppt/slides/slide{i + 1}.xml", xml)

    buf.seek(0)
    return buf.getvalue()


def _xml_escape(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")


async def generate_word(topic: str, grade: GradeLevel) -> bytes:
    """Generate a teaching Word document about a topic, returns .docx bytes."""
    prompt = f"""为「{topic}」生成一份K12教学文档（{grade.value}水平），包含以下部分：

1. 学习目标
2. 核心概念（3-5个）
3. 详细讲解（500-800字）
4. 生活例子
5. 思考题

返回JSON：
{{
  "title": "文档标题",
  "objectives": ["目标1", "目标2"],
  "concepts": [{{"term": "概念名", "explanation": "解释"}}],
  "content": "详细讲解内容（Markdown格式）",
  "example": "生活中的例子",
  "questions": ["思考题1", "思考题2"]
}}"""

    resp = await client.chat.completions.create(
        model=LLM_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=2000,
        response_format={"type": "json_object"},
    )

    try:
        data = json.loads(resp.choices[0].message.content)
    except json.JSONDecodeError:
        data = {"title": topic, "objectives": [], "concepts": [], "content": "", "example": "", "questions": []}

    if not DOCX_AVAILABLE:
        return _make_simple_docx(data)

    return _make_docx(data)


def _make_docx(data: dict) -> bytes:
    doc = Document()
    doc.styles["Normal"].font.size = DocPt(11)
    doc.styles["Normal"].font.name = "Microsoft YaHei"

    doc.add_heading(data.get("title", "教学文档"), level=0)

    objectives = data.get("objectives", [])
    if objectives:
        doc.add_heading("学习目标", level=1)
        for obj in objectives:
            doc.add_paragraph(obj, style="List Bullet")

    concepts = data.get("concepts", [])
    if concepts:
        doc.add_heading("核心概念", level=1)
        for c in concepts:
            doc.add_paragraph(f"{c.get('term', '')}: {c.get('explanation', '')}")

    content = data.get("content", "")
    if content:
        doc.add_heading("详细讲解", level=1)
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            else:
                doc.add_paragraph(line)

    example = data.get("example", "")
    if example:
        doc.add_heading("生活例子", level=1)
        doc.add_paragraph(example)

    questions = data.get("questions", [])
    if questions:
        doc.add_heading("思考题", level=1)
        for i, q in enumerate(questions, 1):
            doc.add_paragraph(f"{i}. {q}")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _make_simple_docx(data: dict) -> bytes:
    """Simple DOCX without python-docx dependency."""
    buf = io.BytesIO()
    import zipfile

    title = _xml_escape(data.get("title", "教学文档"))
    objectives = "".join(f"<w:p><w:r><w:t>{_xml_escape(o)}</w:t></w:r></w:p>" for o in data.get("objectives", []))
    concepts = "".join(
        f"<w:p><w:r><w:rPr><w:b/></w:rPr><w:t>{_xml_escape(c.get('term', ''))}</w:t></w:r><w:r><w:t>: {_xml_escape(c.get('explanation', ''))}</w:t></w:r></w:p>"
        for c in data.get("concepts", [])
    )
    content = "".join(f"<w:p><w:r><w:t>{_xml_escape(line.strip())}</w:t></w:r></w:p>" for line in data.get("content", "").split("\n") if line.strip())
    example = f"<w:p><w:r><w:t>{_xml_escape(data.get('example', ''))}</w:t></w:r></w:p>"
    questions = "".join(f"<w:p><w:r><w:t>{i + 1}. {_xml_escape(q)}</w:t></w:r></w:p>" for i, q in enumerate(data.get("questions", [])))

    document_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:rPr><w:b/><w:sz w:val="48"/></w:rPr><w:t>{title}</w:t></w:r></w:p>
    <w:p><w:r><w:rPr><w:b/><w:sz w:val="32"/></w:rPr><w:t>学习目标</w:t></w:r></w:p>{objectives}
    <w:p><w:r><w:rPr><w:b/><w:sz w:val="32"/></w:rPr><w:t>核心概念</w:t></w:r></w:p>{concepts}
    <w:p><w:r><w:rPr><w:b/><w:sz w:val="32"/></w:rPr><w:t>详细讲解</w:t></w:r></w:p>{content}
    <w:p><w:r><w:rPr><w:b/><w:sz w:val="32"/></w:rPr><w:t>生活例子</w:t></w:r></w:p>{example}
    <w:p><w:r><w:rPr><w:b/><w:sz w:val="32"/></w:rPr><w:t>思考题</w:t></w:r></w:p>{questions}
  </w:body>
</w:document>'''

    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>''')
        z.writestr("_rels/.rels", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>''')
        z.writestr("word/document.xml", document_xml)

    buf.seek(0)
    return buf.getvalue()

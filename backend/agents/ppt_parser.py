"""
PPT Deep Parser — A→D 四层解析管线 + 三级降级策略。

A. 文本提取 (slide-level text extraction, python-pptx)
B. 图片OCR  (提取PPT内嵌图片 → 本地tesseract，免费)
C. 结构理解 (heading/bullet/example recognition, LLM)
D. 互动生成 (quizzes, animations, teaching cards, LLM)

LLM只看文字不看图：L1提取文字 + L2图片OCR → 合并 → 一次性发给LLM做B/C/D
"""

import io
import json
import os
import re
import tempfile
import zipfile
import httpx
from openai import AsyncOpenAI

from config import LLM_API_KEY, LLM_BASE_URL, LLM_MODEL

client = AsyncOpenAI(api_key=LLM_API_KEY, base_url=LLM_BASE_URL, timeout=httpx.Timeout(120.0, connect=10.0))

# ─── L2: Image extraction + OCR ───

def _extract_images_from_pptx(file_path: str) -> list[tuple[int, bytes, str]]:
    """Extract embedded images from a PPTX (ZIP format).
    Returns list of (slide_guess, image_bytes, ext).
    """
    images = []
    try:
        with zipfile.ZipFile(file_path, "r") as z:
            names = z.namelist()
            media_files = [n for n in names if n.startswith("ppt/media/") and not n.endswith("/")]
            slide_files = [n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")]

            # Parse slide→image relationships
            slide_image_map: dict[str, list[str]] = {}
            for slide_name in slide_files:
                rels_name = slide_name.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
                if rels_name in names:
                    try:
                        rels_xml = z.read(rels_name).decode("utf-8")
                        import xml.etree.ElementTree as ET
                        root = ET.fromstring(rels_xml)
                        media_refs = []
                        for rel in root:
                            target = rel.get("Target", "")
                            if "media" in target.lower() or any(
                                target.lower().endswith(ext) for ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp")
                            ):
                                media_refs.append("ppt/slides/" + os.path.normpath(os.path.join(
                                    os.path.dirname(slide_name), target)))
                        slide_image_map[slide_name] = media_refs
                    except Exception:
                        pass

            # Extract images and assign to slides
            slide_nums = sorted([int(n.split("slide")[1].split(".")[0]) for n in slide_files])
            slide_prefixes = {}
            for n in slide_nums:
                slide_prefixes[f"ppt/slides/slide{n}.xml"] = n

            for media_name in media_files:
                data = z.read(media_name)
                ext = media_name.rsplit(".", 1)[-1].lower()
                if ext not in ("png", "jpg", "jpeg", "gif", "bmp", "tiff", "webp"):
                    continue

                # Find which slide this image belongs to
                slide_num = 0
                for slide_name, refs in slide_image_map.items():
                    normalized_refs = [r.replace("\\", "/") for r in refs]
                    if media_name in normalized_refs or any(
                        media_name.endswith(r.split("ppt/media/")[-1] if "ppt/media/" in r else r)
                        for r in normalized_refs
                    ):
                        slide_num = slide_prefixes.get(slide_name, 0)
                        break

                if slide_num == 0:
                    # Fallback: assign to most likely slide based on filename order
                    img_index = media_files.index(media_name)
                    ratio = img_index / max(len(media_files) - 1, 1)
                    slide_num = slide_nums[min(int(ratio * (len(slide_nums) - 1)), len(slide_nums) - 1)] if slide_nums else 0

                images.append((slide_num, data, ext))
    except Exception:
        pass

    return images


def _ocr_image(image_bytes: bytes) -> str:
    """Run tesseract OCR on an image. Returns extracted text or empty string."""
    try:
        from PIL import Image
        import pytesseract

        if os.name == "nt":
            for candidate in [
                r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            ]:
                if os.path.exists(candidate):
                    pytesseract.pytesseract.tesseract_cmd = candidate
                    break

        img = Image.open(io.BytesIO(image_bytes))
        if img.mode != "L":
            img = img.convert("L")
        text = pytesseract.image_to_string(img, lang="chi_sim+eng", config="--psm 6")
        return text.strip()
    except ImportError:
        return ""
    except Exception:
        return ""


MULTIMODAL_PROMPT = """描述这张图片的内容。重点关注：
1. 如果是图表（电路图、流程图、架构图、数据图），描述其结构和关键信息
2. 如果是示意图，说明图上标注了什么
3. 如果是公式或代码截图，转写其内容
4. 如果是普通照片，简要描述画面内容
用中文回答，200字以内。"""


async def _multimodal_analyze_image(image_bytes: bytes, image_ext: str = "png") -> str:
    """Send an image to LLM for visual analysis. Used as fallback for diagrams."""
    import base64

    mime_map = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
                "gif": "image/gif", "webp": "image/webp", "bmp": "image/bmp"}
    mime = mime_map.get(image_ext.lower(), "image/png")
    data_url = f"data:{mime};base64," + base64.b64encode(image_bytes).decode()

    try:
        resp = await client.chat.completions.create(
            model=LLM_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": MULTIMODAL_PROMPT},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }],
            max_tokens=400,
            temperature=0.3,
        )
        return resp.choices[0].message.content.strip()
    except Exception:
        # Model may not support vision — silently skip
        return ""


OCR_FALLBACK_THRESHOLD = 10  # chars: below this, treat as diagram and use multimodal


def _ocr_or_multimodal(image_bytes: bytes, ext: str) -> tuple[str, dict]:
    """Hybrid strategy: OCR first, multimodal fallback for diagrams/charts.
    Returns (text, stats).
    """
    stats = {"method": "none", "chars": 0}

    # 1. Try OCR first (free)
    ocr_text = _ocr_image(image_bytes)
    if ocr_text and len(ocr_text) >= OCR_FALLBACK_THRESHOLD:
        stats["method"] = "ocr"
        stats["chars"] = len(ocr_text)
        return ocr_text, stats

    # 2. OCR got nothing or too little — likely a diagram
    # Note: multimodal runs async, handled separately
    stats["method"] = "multimodal_needed"
    stats["chars"] = len(ocr_text)
    stats["ocr_text"] = ocr_text
    return ocr_text, stats


# ─── Combined extraction (L1 + L2 sync + L2.5 async multimodal) ───

def ocr_stats() -> dict:
    """Check if OCR dependencies are available."""
    try:
        import pytesseract
        from PIL import Image
        return {"tesseract_available": True}
    except ImportError:
        return {"tesseract_available": False}


def extract_slides_text(file_path: str, filename: str) -> tuple[list[dict], dict, list[dict]]:
    """Extract text from each slide (L1) + OCR embedded images (L2).
    Identifies images that need multimodal fallback.

    Returns (slides, stats, multimodal_queue)
      multimodal_queue = [{slide_num, image_bytes, ext}, ...]
    """
    slides = []
    img_stats = {"total_images": 0, "ocr_success": 0, "multimodal_needed": 0, "ocr_chars": 0}
    multimodal_queue = []
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if ext not in ("pptx", "ppt"):
        return slides, img_stats, multimodal_queue

    # ─── L1: Text extraction via python-pptx ───
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                                ph = shape.placeholder_format
                                if ph.type == 1:
                                    if not title:
                                        title = t
                            texts.append(t)
            slides.append({
                "slide_num": i,
                "title": title or (texts[0] if texts else ""),
                "text": "\n".join(texts),
                "ocr_text": "",
                "multimodal_text": "",
            })
    except Exception:
        pass

    # ─── L2: Image extraction + OCR ───
    images = _extract_images_from_pptx(file_path)
    img_stats["total_images"] = len(images)

    if images:
        slide_ocr: dict[int, list[str]] = {}
        for slide_num, img_bytes, img_ext in images:
            ocr_text, decision = _ocr_or_multimodal(img_bytes, img_ext)
            if decision["method"] == "ocr":
                slide_ocr.setdefault(slide_num, []).append("[OCR]" + ocr_text)
                img_stats["ocr_success"] += 1
                img_stats["ocr_chars"] += len(ocr_text)
            elif decision["method"] == "multimodal_needed":
                # Queue for async multimodal processing
                multimodal_queue.append({
                    "slide_num": slide_num,
                    "image_bytes": img_bytes,
                    "ext": img_ext,
                })
                img_stats["multimodal_needed"] += 1
                # Still use OCR text if any
                if ocr_text:
                    slide_ocr.setdefault(slide_num, []).append("[OCR片段]" + ocr_text)

        for s in slides:
            sn = s["slide_num"]
            if sn in slide_ocr:
                s["ocr_text"] = "\n".join(slide_ocr[sn])
                s["text"] = s["text"] + "\n[图片文字]\n" + s["ocr_text"]

    return slides, img_stats, multimodal_queue


async def process_multimodal_fallbacks(slides: list[dict], multimodal_queue: list[dict]) -> dict:
    """Run multimodal vision analysis on images that OCR couldn't handle.
    Merges results back into slides.
    Returns stats.
    """
    import asyncio

    stats = {"multimodal_total": len(multimodal_queue), "multimodal_success": 0, "multimodal_chars": 0}

    if not multimodal_queue:
        return stats

    # Process concurrently (max 3 at a time to avoid rate limits)
    semaphore = asyncio.Semaphore(3)

    async def process_one(item: dict):
        async with semaphore:
            text = await _multimodal_analyze_image(item["image_bytes"], item["ext"])
            return item["slide_num"], text

    tasks = [process_one(item) for item in multimodal_queue]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    # Merge results into slides
    slide_mm: dict[int, list[str]] = {}
    for result in results:
        if isinstance(result, Exception):
            continue
        slide_num, text = result
        if text:
            slide_mm.setdefault(slide_num, []).append("[多模态]" + text)
            stats["multimodal_success"] += 1
            stats["multimodal_chars"] += len(text)

    for s in slides:
        sn = s["slide_num"]
        if sn in slide_mm:
            s["multimodal_text"] = "\n".join(slide_mm[sn])
            s["text"] = s["text"] + "\n[图表分析]\n" + s["multimodal_text"]

    return stats


# ─── Step B+C+D: LLM-based deep analysis (same as before) ───

PARSE_PROMPT = """你是一个教育AI专家。分析以下教学PPT的内容，做4层深度解析。

## 输入的PPT内容
{slides_text}

## 要求输出JSON格式

{{
  "course_meta": {{
    "title": "课程标题（从PPT推断）",
    "subject": "所属学科/领域",
    "estimated_grade": "primary_low|primary_high|middle|high",
    "estimated_duration": "预估教学时长（分钟）",
    "difficulty": "beginner|intermediate|advanced",
    "prerequisites": ["前置知识1", "前置知识2"]
  }},
  "slide_analysis": [
    {{
      "slide_num": 1,
      "type": "title|content|example|exercise|summary|transition",
      "heading": "本页标题",
      "key_message": "本页核心信息（一句话）",
      "knowledge_points": ["知识点1", "知识点2"],
      "importance": "core|supporting|optional"
    }}
  ],
  "knowledge_graph": [
    {{"source": "知识点A", "relation": "前置于|包含|关联", "target": "知识点B"}}
  ],
  "key_concepts": [
    {{
      "term": "核心概念名",
      "definition": "简洁定义（50字内）",
      "grade_level": "适合哪个学段",
      "examples": ["例子1", "例子2"]
    }}
  ],
  "generated_quiz": [
    {{
      "type": "choice",
      "question": "题目",
      "options": ["A. x", "B. x", "C. x", "D. x"],
      "answer": "A",
      "explanation": "解析",
      "related_knowledge": "考察的知识点"
    }}
  ],
  "teaching_suggestions": {{
    "animation_topics": ["适合做动画的概念"],
    "coding_examples": ["适合编程实践的主题"],
    "discussion_questions": ["课堂讨论题"],
    "recommended_next": "学完这个PPT后建议学什么"
  }},
  "knowledge_entries": [
    {{
      "grade": "primary_low|primary_high|middle|high",
      "topic": "所属主题",
      "title": "知识条目标题",
      "content": "300-500字的知识讲解，直接可以加入RAG知识库"
    }}
  ]
}}

## 注意事项
- 知识点之间要标注关联关系（knowledge_graph）
- 生成的题目要覆盖PPT的核心知识点
- knowledge_entries要有实质内容，每个条目200-500字，可以直接入库
- 如果PPT某页是纯标题/过渡页，type标记为transition
- 尽量多地提取和生成内容，不要偷懒"""


async def deep_parse_ppt(slides: list[dict], filename: str) -> dict:
    """Send all slide content to LLM for deep 4-layer analysis."""
    # Build slides text representation
    slides_text_parts = []
    for s in slides:
        slides_text_parts.append(
            f"--- 第{s['slide_num']}页 ---\n"
            f"标题: {s.get('title', '无')}\n"
            f"内容: {s.get('text', '无')[:800]}\n"
        )
    slides_text = "\n".join(slides_text_parts)

    # Truncate if too long (reserve ~60K chars for LLM context)
    if len(slides_text) > 30000:
        slides_text = slides_text[:30000] + "\n\n[内容太长已截断，请基于前30页分析]"

    try:
        resp = await client.chat.completions.create(
            model=LLM_MODEL,
            messages=[{"role": "user", "content": PARSE_PROMPT.format(slides_text=slides_text)}],
            temperature=0.5,
            max_tokens=8000,
            response_format={"type": "json_object"},
        )
        raw = resp.choices[0].message.content

        # Clean potential markdown code fences
        if "```json" in raw:
            raw = raw.split("```json")[1].split("```")[0]
        elif "```" in raw:
            raw = raw.split("```")[1].split("```")[0]

        return json.loads(raw)
    except (json.JSONDecodeError, Exception) as e:
        return {
            "error": str(e),
            "course_meta": {"title": filename, "subject": "未知", "estimated_grade": "middle", "difficulty": "intermediate"},
            "slide_analysis": [],
            "knowledge_graph": [],
            "key_concepts": [],
            "generated_quiz": [],
            "teaching_suggestions": {},
            "knowledge_entries": [],
        }


# ─── Post-processing: index generated knowledge into RAG ───

async def index_parsed_knowledge(knowledge_entries: list[dict]):
    """Add LLM-generated knowledge entries into the ChromaDB RAG index."""
    if not knowledge_entries:
        return 0

    try:
        from utils.rag import rag
        docs = [e.get("content", "") for e in knowledge_entries]
        metas = [
            {
                "grade": e.get("grade", "middle"),
                "topic": e.get("topic", "parsed"),
                "title": e.get("title", ""),
            }
            for e in knowledge_entries
        ]
        ids = [f"ppt_parsed_{i}_{hash(d) % 100000}" for i, d in enumerate(docs)]
        embeddings = rag.embed_model.encode(docs).tolist()
        rag.collection.add(embeddings=embeddings, documents=docs, metadatas=metas, ids=ids)
        return len(docs)
    except Exception:
        return 0

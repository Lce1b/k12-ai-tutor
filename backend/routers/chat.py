"""
Chat API router — handles all user interactions.
"""

import logging
import time
from fastapi import APIRouter
from fastapi.responses import Response
from pydantic import BaseModel

from agents.orchestrator import (
    classify_intent,
    extract_topic,
    generate_greeting,
    suggest_next_step,
    GradeLevel,
    Intent,
)
from agents.dialogue import chat as dialogue_chat
from agents.teaching import teach
from agents.animation import generate_animation
from agents.quiz import generate_quiz, evaluate_answer
from agents.coding import generate_code_example, execute_code, get_code_templates, get_template_by_id
from agents.picturebook import generate_story
from agents.resources import generate_ppt, generate_word
from agents.ppt_parser import extract_slides_text, process_multimodal_fallbacks, deep_parse_ppt, index_parsed_knowledge, ocr_stats
from db import execute, fetchall, fetchone
from utils.rag import rag
from utils.history import (
    record_message, record_quiz_result, get_stats, add_xp, record_streak,
    ensure_session,
)

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api", tags=["chat"])


class ChatRequest(BaseModel):
    session_id: str = "default"
    message: str
    grade: str = "middle"


class QuizEvalRequest(BaseModel):
    question: dict
    answer: str
    session_id: str = "default"
    topic: str = ""


class CodeExecRequest(BaseModel):
    code: str


class GreetingRequest(BaseModel):
    session_id: str = "default"
    grade: str = "middle"


class ResourceRequest(BaseModel):
    topic: str
    grade: str = "middle"
    session_id: str = "default"


@router.get("/greeting")
async def greeting(session_id: str = "default", grade: str = "middle"):
    """Generate a proactive greeting with learning suggestions."""
    g = GradeLevel(grade)
    greeting_data = await generate_greeting(g)
    return {
        **greeting_data,
        "grade": grade,
        "stats": await get_stats(session_id),
    }


@router.post("/chat")
async def handle_chat(req: ChatRequest):
    """Main chat endpoint — orchestrates between agents."""
    grade = GradeLevel(req.grade)
    intent = await classify_intent(req.message)
    topic = extract_topic(req.message)

    # Load recent history from MySQL
    db_history = await fetchall(
        "SELECT role, content FROM messages WHERE session_id = %s ORDER BY created_at DESC LIMIT 20",
        (req.session_id,),
    )
    history = [{"role": r["role"], "content": r["content"]} for r in reversed(db_history)]

    rag_results = rag.search(req.message, grade=req.grade, top_k=3)
    rag_context = "\n---\n".join(r["content"][:300] for r in rag_results)

    result = {
        "intent": intent.value,
        "grade": grade.value,
        "topic": topic,
        "rag_sources": [{"title": r["meta"].get("title", ""), "score": round(r["score"], 3)} for r in rag_results],
    }

    try:
        if intent == Intent.TEACH:
            lesson = await teach(topic, history, grade)
            result["type"] = "lesson"
            result["lesson"] = lesson
            result["message"] = f"## {lesson.get('title', topic)}\n\n{lesson.get('intro', '')}"

        elif intent == Intent.ANIMATE:
            html = await generate_animation(topic, grade)
            result["type"] = "animation"
            result["animation_html"] = html
            result["message"] = f"我为你生成了一个讲解「{topic}」的动画，请在右侧查看~"

        elif intent == Intent.QUIZ:
            quiz = await generate_quiz(topic, grade)
            result["type"] = "quiz"
            result["quiz"] = quiz
            qs = quiz.get("questions", [])
            result["message"] = f"来挑战一下关于「{topic}」的练习吧！共{len(qs)}道题。"

        elif intent == Intent.CODE:
            code_data = await generate_code_example(topic, grade)
            result["type"] = "code"
            result["code"] = code_data
            result["message"] = f"## {code_data.get('title', '')}\n\n{code_data.get('explanation', '')}\n\n```python\n{code_data.get('code', '')}\n```"

        elif intent == Intent.PICTURE_BOOK:
            story = await generate_story(topic, grade)
            result["type"] = "picture_book"
            result["story"] = story
            result["message"] = f"## {story.get('title', '')}\n\n" + "\n\n".join(
                p["text"] for p in story.get("pages", [])
            )

        else:  # CHAT
            reply = await dialogue_chat(req.message, history, grade, intent, rag_context)
            result["type"] = "chat"
            result["message"] = reply

    except Exception:
        logger.exception("Agent dispatch failed, falling back to chat")
        reply = await dialogue_chat(req.message, history, grade, Intent.CHAT, rag_context)
        result["type"] = "chat"
        result["message"] = reply

    # Generate proactive next-step suggestion
    try:
        next_step = await suggest_next_step(topic, intent, grade)
        result["next_step"] = next_step
    except Exception:
        result["next_step"] = {"message": "学得不错！还想了解什么？", "suggestions": []}

    # Persist messages to MySQL
    await ensure_session(req.session_id, req.grade)
    await execute(
        "INSERT INTO messages (session_id, role, content, intent, topic) VALUES (%s, %s, %s, %s, %s)",
        (req.session_id, "user", req.message, intent.value, topic),
    )
    await execute(
        "INSERT INTO messages (session_id, role, content, intent, topic) VALUES (%s, %s, %s, %s, %s)",
        (req.session_id, "assistant", result["message"][:500], intent.value, topic),
    )

    await record_message(req.session_id, req.grade, topic, intent.value)

    return result


@router.post("/quiz/eval")
async def quiz_eval(req: QuizEvalRequest):
    result = await evaluate_answer(req.question, req.answer)
    await record_quiz_result(req.session_id, req.topic, result["is_correct"])

    # Gamification: award XP and track streak
    xp_amount = 10 if result["is_correct"] else 2  # partial XP even on wrong answers
    xp_info = await add_xp(req.session_id, xp_amount)
    streak_info = await record_streak(req.session_id, result["is_correct"])

    return {
        **result,
        "xp": xp_info,
        "streak": streak_info,
    }


@router.post("/code/exec")
async def code_exec(req: CodeExecRequest):
    return await execute_code(req.code)


@router.get("/session/{session_id}")
async def get_session_history(session_id: str):
    rows = await fetchall(
        "SELECT role, content, intent, topic, created_at FROM messages WHERE session_id = %s ORDER BY created_at ASC LIMIT 50",
        (session_id,),
    )
    return {
        "history": [dict(r) for r in rows],
        "stats": await get_stats(session_id),
    }


@router.delete("/session/{session_id}")
async def clear_session(session_id: str):
    await execute("DELETE FROM messages WHERE session_id = %s", (session_id,))
    return {"status": "ok"}


@router.get("/stats/{session_id}")
async def learning_stats(session_id: str):
    return await get_stats(session_id)


# ---- Curriculum endpoint ----

import json
import os

@router.get("/curriculum")
async def get_curriculum():
    """Return the K12 AI course curriculum."""
    path = os.path.join(os.path.dirname(__file__), "..", "knowledge", "curriculum.json")
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


# ---- Code templates ----

@router.get("/code/templates")
async def code_templates(grade: str = "middle"):
    return {"templates": get_code_templates(grade)}


@router.get("/code/templates/{template_id}")
async def code_template_detail(template_id: str):
    tmpl = get_template_by_id(template_id)
    if tmpl is None:
        from fastapi.responses import JSONResponse
        return JSONResponse({"error": "Template not found"}, status_code=404)
    return tmpl


# ---- File upload for teaching resources ----

import shutil
from fastapi import UploadFile, File, Form

UPLOAD_DIR = os.path.join(os.path.dirname(__file__), "..", "uploads")

@router.post("/resources/upload")
async def upload_resource(
    file: UploadFile = File(...),
    session_id: str = Form("default"),
    grade: str = Form("middle"),
):
    """Upload a teaching resource (PDF, DOC, PPT, image, video)."""
    os.makedirs(UPLOAD_DIR, exist_ok=True)

    # Save file
    safe_name = f"{int(time.time())}_{file.filename}"
    file_path = os.path.join(UPLOAD_DIR, safe_name)
    with open(file_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    # Extract text from common formats for indexing
    text_content = _extract_text(file_path, file.filename)

    # Index into knowledge base if text extracted
    if text_content:
        try:
            from utils.rag import rag
            rag.collection.add(
                embeddings=rag.embed_model.encode([text_content[:500]]).tolist(),
                documents=[text_content[:500]],
                metadatas=[{"grade": grade, "topic": "uploaded", "title": file.filename}],
                ids=[f"upload_{safe_name}"],
            )
        except Exception:
            pass

    await record_message(session_id, grade, file.filename, "resource_upload")

    return {
        "filename": safe_name,
        "original_name": file.filename,
        "size": os.path.getsize(file_path),
        "text_preview": text_content[:200] if text_content else "",
    }


def _extract_text(path: str, filename: str) -> str:
    """Try to extract text from uploaded file."""
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if ext == "txt":
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()[:2000]

    if ext == "pdf":
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            return "".join((p.extract_text() or "") for p in reader.pages)[:2000]
        except Exception:
            pass

    if ext in ("docx", "doc"):
        try:
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)[:2000]
        except Exception:
            pass

    if ext in ("pptx", "ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)[:2000]
        except Exception:
            pass

    return ""


@router.get("/resources/ocr-status")
async def ocr_status():
    """Check if OCR dependencies (tesseract + pillow) are available."""
    return ocr_stats()


# ---- PPT Deep Parse (A→D pipeline) ----

@router.post("/resources/ppt/deep-parse")
async def deep_parse_ppt_endpoint(
    file: UploadFile = File(...),
    grade: str = Form("middle"),
    auto_index: bool = Form(True),
):
    """Upload a PPT and run the full A→D 4-layer analysis pipeline.

    A. Extract text per slide
    B. Understand slide structure (heading/content/example/exercise)
    C. Understand teaching content (topic, difficulty, knowledge graph)
    D. Generate quizzes, animations, teaching cards

    Returns the full analysis result. If auto_index=True, generated
    knowledge entries are automatically added to the RAG index.
    """
    # Save uploaded file
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    safe_name = f"ppt_{int(time.time())}_{file.filename}"
    file_path = os.path.join(UPLOAD_DIR, safe_name)
    with open(file_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    # L1: Extract text + L2: OCR images (sync, free)
    slides, img_stats, multimodal_queue = extract_slides_text(file_path, file.filename)
    if not slides:
        return {"error": "无法提取PPT内容，请确认文件格式正确（.pptx）", "slides": []}

    # L2.5: Multimodal fallback for diagrams/charts (async, costs LLM tokens)
    mm_stats = await process_multimodal_fallbacks(slides, multimodal_queue)

    # L3-5: Deep LLM analysis (LLM only sees merged text, never raw images)
    result = await deep_parse_ppt(slides, file.filename)

    if "error" not in result:
        result["slides_raw"] = [{
            "slide_num": s["slide_num"],
            "title": s.get("title", ""),
            "text_len": len(s.get("text", "")),
            "has_ocr": bool(s.get("ocr_text", "")),
            "has_multimodal": bool(s.get("multimodal_text", "")),
        } for s in slides]
        result["total_slides"] = len(slides)
        result["filename"] = file.filename
        result["image_stats"] = {
            "total_images": img_stats["total_images"],
            "ocr_handled": img_stats["ocr_success"],
            "ocr_chars": img_stats["ocr_chars"],
            "multimodal_handled": mm_stats["multimodal_success"],
            "multimodal_total": mm_stats["multimodal_total"],
            "multimodal_chars": mm_stats["multimodal_chars"],
        }

        # Auto-index generated knowledge entries
        if auto_index and result.get("knowledge_entries"):
            indexed = await index_parsed_knowledge(result["knowledge_entries"])
            result["indexed_count"] = indexed
        else:
            result["indexed_count"] = 0

    return result


# ---- Resource generation endpoints ----

@router.post("/resources/ppt")
async def download_ppt(req: ResourceRequest):
    """Generate and download a teaching PPT."""
    g = GradeLevel(req.grade)
    pptx_bytes = await generate_ppt(req.topic, g)
    await record_message(req.session_id, req.grade, req.topic, "resource_ppt")
    filename = f"{req.topic}_教学课件.pptx"
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{filename}"},
    )


@router.post("/resources/word")
async def download_word(req: ResourceRequest):
    """Generate and download a teaching Word document."""
    g = GradeLevel(req.grade)
    docx_bytes = await generate_word(req.topic, g)
    await record_message(req.session_id, req.grade, req.topic, "resource_word")
    filename = f"{req.topic}_教学文档.docx"
    return Response(
        content=docx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{filename}"},
    )
